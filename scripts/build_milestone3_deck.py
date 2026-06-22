from __future__ import annotations

import csv
import json
import subprocess
from pathlib import Path

import imageio_ffmpeg
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[1]
DOCS = REPO / "docs"
ASSETS = DOCS / "assets" / "milestone3_preliminary_results"
BENCH = (
    REPO
    / "data"
    / "benchmarks"
    / "distilled098_24fps_49ctx_72future_base_vs_full_lora_lr5e6_step3000_seed231_all5"
)
RANK_SWEEP = (
    REPO
    / "data"
    / "benchmarks"
    / "distilled098_24fps_49ctx_72future_rank_sweep_lr5e6_step3000_seed231_all5"
    / "rank_sweep_summary.csv"
)
CHECKPOINT_SWEEP = (
    REPO
    / "data"
    / "benchmarks"
    / "distilled098_checkpoint_sweep_24fps_49ctx_72future_seed231_all_ranks"
    / "checkpoint_sweep_summary_with_quality_gate.csv"
)
MANIFEST = (
    REPO
    / "data"
    / "distilled098_24fps_49ctx_72future_base_vs_full_lora_lr5e6_step3000_seed231_all5"
    / "manifest_24fps_49ctx_72future_base_vs_full_lora_lr5e6_step3000_seed231_all5.json"
)
SOURCE_DIR = REPO / "data" / "inference_input_clips" / "interpolated_24fps_waymo_full20s"
OUT_PPTX = DOCS / "cs231n_milestone3_preliminary_results_may29.pptx"


BG = RGBColor(247, 245, 239)
INK = RGBColor(26, 32, 44)
MUTED = RGBColor(88, 95, 105)
BLUE = RGBColor(29, 78, 216)
TEAL = RGBColor(15, 118, 110)
RED = RGBColor(185, 28, 28)
AMBER = RGBColor(180, 83, 9)
WHITE = RGBColor(255, 255, 255)


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def as_float(row: dict[str, str], key: str) -> float:
    return float(row[key])


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(31)
    run.font.bold = True
    run.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.02), Inches(11.8), Inches(0.36))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Aptos"
        sr.font.size = Pt(13)
        sr.font.color.rgb = MUTED


def add_footer(slide, n: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_callout(slide, text: str, x: float, y: float, w: float, h: float, color=BLUE) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos Display"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_metric_card(slide, label: str, value: str, x: float, y: float, color=BLUE, note: str = "") -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.55), Inches(1.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 216, 206)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = "Aptos"
    p2.font.size = Pt(10)
    p2.font.color.rgb = INK
    if note:
        p3 = tf.add_paragraph()
        p3.text = note
        p3.font.name = "Aptos"
        p3.font.size = Pt(8)
        p3.font.color.rgb = MUTED


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c in range(len(rows[0])):
        table.columns[c].width = Inches(w / len(rows[0]))
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(235, 231, 222) if r_idx == 0 else WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Aptos"
            para.font.size = Pt(10 if r_idx else 9)
            para.font.bold = r_idx == 0
            para.font.color.rgb = INK


def make_charts(summary_rows: list[dict[str, str]], report: dict) -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    labels = ["Base\nno LoRA", "Full LoRA\nstep 3000"]
    base, lora = summary_rows

    plt.style.use("default")
    fig, ax = plt.subplots(figsize=(7.8, 4.2), dpi=180)
    x = [0, 1]
    psnr = [as_float(base, "mean_future_psnr"), as_float(lora, "mean_future_psnr")]
    ssim = [as_float(base, "mean_future_global_ssim"), as_float(lora, "mean_future_global_ssim")]
    bars = ax.bar([v - 0.18 for v in x], psnr, width=0.34, label="Future PSNR ↑", color="#1d4ed8")
    ax2 = ax.twinx()
    ax2.bar([v + 0.18 for v in x], ssim, width=0.34, label="Future SSIM ↑", color="#0f766e")
    ax.set_xticks(x, labels)
    ax.set_ylim(0, max(psnr) * 1.25)
    ax2.set_ylim(0.70, 0.77)
    ax.set_ylabel("PSNR")
    ax2.set_ylabel("Global SSIM")
    ax.set_title("Reference Metrics: Full-Data LoRA Improves Future Similarity")
    for bar, value in zip(bars, psnr):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.25, f"{value:.2f}", ha="center", fontsize=9)
    for xpos, value in zip([v + 0.18 for v in x], ssim):
        ax2.text(xpos, value + 0.002, f"{value:.3f}", ha="center", fontsize=9)
    fig.tight_layout()
    ref_chart = ASSETS / "reference_metrics.png"
    fig.savefig(ref_chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(7.8, 4.2), dpi=180)
    sharp = [
        as_float(base, "mean_sharpness_ratio_generated_over_reference"),
        as_float(lora, "mean_sharpness_ratio_generated_over_reference"),
    ]
    motion = [
        as_float(base, "mean_motion_ratio_generated_over_reference"),
        as_float(lora, "mean_motion_ratio_generated_over_reference"),
    ]
    ax.bar([v - 0.18 for v in x], sharp, width=0.34, label="Sharpness ratio ↑", color="#b45309")
    ax.bar([v + 0.18 for v in x], motion, width=0.34, label="Motion ratio ↑", color="#be123c")
    ax.axhline(1.0, color="#374151", linestyle="--", linewidth=1, label="Real future = 1.0")
    ax.set_xticks(x, labels)
    ax.set_ylim(0, 1.1)
    ax.set_title("Quality Diagnostics: Motion Recovers, Sharpness Still Lags")
    ax.legend(loc="upper right", fontsize=8)
    for xpos, value in zip([v - 0.18 for v in x], sharp):
        ax.text(xpos, value + 0.025, f"{value:.3f}", ha="center", fontsize=9)
    for xpos, value in zip([v + 0.18 for v in x], motion):
        ax.text(xpos, value + 0.025, f"{value:.3f}", ha="center", fontsize=9)
    fig.tight_layout()
    quality_chart = ASSETS / "quality_diagnostics.png"
    fig.savefig(quality_chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    gate = report["quality_gate"]["observed"]
    fig, ax = plt.subplots(figsize=(7.8, 4.2), dpi=180)
    gate_labels = ["Sharpness\nretention", "Motion\nretention", "SSIM delta\n+ 1"]
    gate_values = [
        gate["sharpness_retention_lora_over_base"],
        gate["motion_retention_lora_over_base"],
        gate["future_global_ssim_delta_lora_minus_base"] + 1.0,
    ]
    colors = [
        "#b91c1c" if gate_values[0] < 0.8 else "#0f766e",
        "#b91c1c" if gate_values[1] < 0.8 else "#0f766e",
        "#b91c1c" if gate_values[2] < 0.99 else "#0f766e",
    ]
    ax.bar(gate_labels, gate_values, color=colors)
    ax.axhline(0.8, color="#111827", linestyle="--", linewidth=1, label="Retention gate = 0.8")
    ax.axhline(0.99, color="#0f766e", linestyle=":", linewidth=1, label="SSIM gate = 0.99")
    ax.set_ylim(0, 1.12)
    ax.set_title("Quality Gate: Narrow Failure Due To Sharpness")
    ax.legend(fontsize=8, loc="lower right")
    for xpos, value in enumerate(gate_values):
        ax.text(xpos, value + 0.025, f"{value:.3f}", ha="center", fontsize=9)
    fig.tight_layout()
    gate_chart = ASSETS / "quality_gate.png"
    fig.savefig(gate_chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    return {"reference": ref_chart, "quality": quality_chart, "gate": gate_chart}


def make_rank_sweep_chart(rank_rows: list[dict[str, str]]) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    rank_rows = sorted(rank_rows, key=lambda r: int(r["rank"]))
    labels = [f"r{row['rank']}" for row in rank_rows]
    ssim = [as_float(row, "mean_future_global_ssim") for row in rank_rows]
    psnr = [as_float(row, "mean_future_psnr") for row in rank_rows]
    sharp = [as_float(row, "mean_sharpness_ratio_generated_over_reference") for row in rank_rows]
    motion = [as_float(row, "mean_motion_ratio_generated_over_reference") for row in rank_rows]

    plt.style.use("default")
    fig, axes = plt.subplots(1, 2, figsize=(9.2, 3.8), dpi=180)
    x = range(len(labels))
    axes[0].bar([v - 0.18 for v in x], psnr, width=0.34, color="#1d4ed8", label="PSNR ↑")
    ax2 = axes[0].twinx()
    ax2.bar([v + 0.18 for v in x], ssim, width=0.34, color="#0f766e", label="SSIM ↑")
    axes[0].set_xticks(list(x), labels)
    axes[0].set_title("Future Similarity")
    axes[0].set_ylabel("PSNR")
    ax2.set_ylabel("SSIM")
    axes[0].set_ylim(16.7, 18.0)
    ax2.set_ylim(0.68, 0.78)
    for xpos, value in zip([v - 0.18 for v in x], psnr):
        axes[0].text(xpos, value + 0.03, f"{value:.2f}", ha="center", fontsize=8)
    for xpos, value in zip([v + 0.18 for v in x], ssim):
        ax2.text(xpos, value + 0.002, f"{value:.3f}", ha="center", fontsize=8)

    axes[1].bar([v - 0.18 for v in x], sharp, width=0.34, color="#b45309", label="Sharpness ratio ↑")
    axes[1].bar([v + 0.18 for v in x], motion, width=0.34, color="#be123c", label="Motion ratio ↑")
    axes[1].axhline(0.8, color="#111827", linestyle="--", linewidth=1, label="Retention gate")
    axes[1].set_xticks(list(x), labels)
    axes[1].set_ylim(0, 1.0)
    axes[1].set_title("Quality Diagnostics")
    axes[1].legend(fontsize=7, loc="upper right")
    for xpos, value in zip([v - 0.18 for v in x], sharp):
        axes[1].text(xpos, value + 0.025, f"{value:.3f}", ha="center", fontsize=8)
    for xpos, value in zip([v + 0.18 for v in x], motion):
        axes[1].text(xpos, value + 0.025, f"{value:.3f}", ha="center", fontsize=8)

    fig.suptitle("Narrow LoRA Rank Sweep: r16 Is The Best Tradeoff So Far", fontsize=12)
    fig.tight_layout()
    chart = ASSETS / "rank_sweep_metrics.png"
    fig.savefig(chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return chart


def summarize_checkpoint_sweep(rows: list[dict[str, str]]) -> dict[str, object]:
    parsed: list[dict[str, object]] = []
    for row in rows:
        parsed.append(
            {
                "rank": int(row["rank"]),
                "step": int(row["step"]),
                "future_psnr": float(row["future_psnr"]),
                "future_ssim": float(row["future_ssim"]),
                "future_mse": float(row["future_mse"]),
                "sharpness_ratio": float(row["sharpness_ratio"]),
                "motion_ratio": float(row["motion_ratio"]),
                "ssim_delta_vs_base": float(row["ssim_delta_vs_base"]),
                "sharpness_retention_vs_base": float(row["sharpness_retention_vs_base"]),
                "motion_retention_vs_base": float(row["motion_retention_vs_base"]),
                "quality_gate": row["quality_gate"],
            }
        )
    passing = [row for row in parsed if row["quality_gate"] == "pass"]
    best_gate = max(passing, key=lambda row: (row["future_ssim"], row["future_psnr"]))
    raw_best = max(parsed, key=lambda row: (row["future_ssim"], row["future_psnr"]))
    best_by_rank: dict[int, dict[str, object]] = {}
    for rank in sorted({int(row["rank"]) for row in parsed}):
        rank_rows = [row for row in parsed if row["rank"] == rank]
        rank_passing = [row for row in rank_rows if row["quality_gate"] == "pass"]
        best_by_rank[rank] = max(rank_passing or rank_rows, key=lambda row: (row["future_ssim"], row["future_psnr"]))
    return {
        "rows": parsed,
        "best_gate": best_gate,
        "raw_best": raw_best,
        "best_by_rank": best_by_rank,
        "num_pass": len(passing),
        "num_total": len(parsed),
    }


def fmt_rank_step(row: dict[str, object]) -> str:
    return f"r{int(row['rank'])} / {int(row['step'])}"


def make_checkpoint_sweep_chart(rows: list[dict[str, str]], base: dict[str, str]) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    parsed = summarize_checkpoint_sweep(rows)
    typed_rows = parsed["rows"]
    best_gate = parsed["best_gate"]
    raw_best = parsed["raw_best"]
    base_ssim = as_float(base, "mean_future_global_ssim")

    plt.style.use("default")
    fig, axes = plt.subplots(1, 2, figsize=(10.4, 3.8), dpi=180)

    for rank, color in [(8, "#0f766e"), (16, "#1d4ed8"), (32, "#b45309")]:
        rank_rows = sorted([row for row in typed_rows if row["rank"] == rank], key=lambda row: row["step"])
        steps = [row["step"] for row in rank_rows]
        ssim = [row["future_ssim"] for row in rank_rows]
        axes[0].plot(steps, ssim, marker="o", linewidth=1.8, color=color, label=f"rank {rank}")
        passing = [row for row in rank_rows if row["quality_gate"] == "pass"]
        if passing:
            axes[0].scatter(
                [row["step"] for row in passing],
                [row["future_ssim"] for row in passing],
                s=72,
                facecolors="none",
                edgecolors="#16a34a",
                linewidths=1.8,
                zorder=3,
            )
    axes[0].axhline(base_ssim, color="#111827", linestyle="--", linewidth=1, label="base SSIM")
    axes[0].set_title("Future SSIM Across Checkpoints")
    axes[0].set_xlabel("training step")
    axes[0].set_ylabel("future SSIM")
    axes[0].set_ylim(0.68, 0.77)
    axes[0].legend(fontsize=7, loc="lower right")

    labels = [f"best valid\n{fmt_rank_step(best_gate)}", f"raw best\n{fmt_rank_step(raw_best)}"]
    sharp = [best_gate["sharpness_retention_vs_base"], raw_best["sharpness_retention_vs_base"]]
    motion = [best_gate["motion_retention_vs_base"], raw_best["motion_retention_vs_base"]]
    x = range(len(labels))
    axes[1].bar([v - 0.18 for v in x], sharp, width=0.34, color="#b45309", label="sharpness retention")
    axes[1].bar([v + 0.18 for v in x], motion, width=0.34, color="#be123c", label="motion retention")
    axes[1].axhline(0.8, color="#111827", linestyle="--", linewidth=1, label="gate")
    axes[1].set_xticks(list(x), labels)
    axes[1].set_ylim(0, 1.12)
    axes[1].set_title("Gate Tradeoff")
    axes[1].legend(fontsize=7, loc="lower right")
    for xpos, value in zip([v - 0.18 for v in x], sharp):
        axes[1].text(xpos, value + 0.025, f"{value:.2f}", ha="center", fontsize=8)
    for xpos, value in zip([v + 0.18 for v in x], motion):
        axes[1].text(xpos, value + 0.025, f"{value:.2f}", ha="center", fontsize=8)

    fig.suptitle("Full Checkpoint Sweep: Best Valid Checkpoint Is r8 Step 1500", fontsize=12)
    fig.tight_layout()
    chart = ASSETS / "checkpoint_sweep_metrics.png"
    fig.savefig(chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return chart


def ffmpeg_extract_frame(video: Path, output: Path, frame_index: int = 80) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
    cmd = [
        ffmpeg,
        "-hide_banner",
        "-loglevel",
        "error",
        "-i",
        str(video),
        "-vf",
        f"fps=24,scale=512:512:force_original_aspect_ratio=increase,crop=512:512,select=eq(n\\,{frame_index})",
        "-frames:v",
        "1",
        "-y",
        str(output),
    ]
    subprocess.run(cmd, check=True)


def make_qual_frames(manifest: dict) -> list[tuple[str, Path]]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    scene = "0081c4821701"
    records = [r for r in manifest["records"] if r["scene_token"] == scene]
    base = next(r for r in records if r["model_mode"] == "base_distilled_no_lora")
    lora = next(r for r in records if r["model_mode"] != "base_distilled_no_lora")
    source = SOURCE_DIR / base["source_filename"]
    outputs = [
        ("Real future", source, ASSETS / "qual_real_future_frame80.png"),
        ("Base distilled", REPO / base["local_file"], ASSETS / "qual_base_frame80.png"),
        ("Full LoRA step 3000", REPO / lora["local_file"], ASSETS / "qual_lora_frame80.png"),
    ]
    result: list[tuple[str, Path]] = []
    for label, video, out in outputs:
        ffmpeg_extract_frame(video, out)
        result.append((label, out))
    return result


def build_deck() -> None:
    summary_rows = read_csv(BENCH / "model_summary.csv")
    summary_rows = sorted(summary_rows, key=lambda r: 0 if r["model_mode"] == "base_distilled_no_lora" else 1)
    rank_rows = read_csv(RANK_SWEEP)
    checkpoint_rows = read_csv(CHECKPOINT_SWEEP)
    report = json.loads((BENCH / "benchmark_report.json").read_text(encoding="utf-8"))
    manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
    base, lora = summary_rows
    checkpoint_summary = summarize_checkpoint_sweep(checkpoint_rows)
    best_gate = checkpoint_summary["best_gate"]
    raw_best = checkpoint_summary["raw_best"]
    best_by_rank = checkpoint_summary["best_by_rank"]
    charts = make_charts(summary_rows, report)
    rank_chart = make_rank_sweep_chart(rank_rows)
    checkpoint_chart = make_checkpoint_sweep_chart(checkpoint_rows, base)
    qual_frames = make_qual_frames(manifest)

    gate = report["quality_gate"]["observed"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []

    def new_slide(title: str, subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, title, subtitle)
        slides.append(slide)
        return slide

    slide = new_slide("Milestone 3 — Preliminary Results", "Action-Conditioned Video World Models for Long-Tail Driving")
    add_bullets(
        slide,
        ["Stanford CS231N · May 29, 2026", "Andrew Liang and Maleeka Raddygala"],
        0.7,
        2.0,
        6.5,
        1.2,
        size=18,
    )
    add_callout(
        slide,
        "Main result: checkpoint selection matters. r8 step 1500 is the best gate-passing checkpoint; r16 step 3000 wins raw SSIM but over-smooths.",
        0.75,
        4.25,
        11.8,
        0.8,
        BLUE,
    )

    slide = new_slide("Question For This Milestone", "Before action conditioning, can we adapt the visual prior without damaging it?")
    add_bullets(
        slide,
        [
            "Long-term objective: generate plausible future driving video conditioned on ego actions.",
            "Current prerequisite: visual continuation on Waymo-style front-camera clips.",
            "We compare base LTX-2B distilled against a corrected full-data visual LoRA checkpoint.",
            "Success means LoRA improves Waymo alignment without losing sharpness, motion, or scene stability.",
        ],
        0.8,
        1.55,
        11.2,
        3.8,
    )

    slide = new_slide("What We Built Since Milestone 2")
    add_bullets(
        slide,
        [
            "24 FPS Waymo front-camera inference set from local 20 second clips.",
            "Full-data LoRA training on 7,992 cached latent windows using LR=5e-6.",
            "Full checkpoint sweep across ranks 8, 16, and 32 at all saved checkpoints.",
            "Fixed-seed inference for base distilled and full-data distilled LoRA checkpoints.",
            "Side-by-side qualitative comparisons across FPS and model variants.",
            "A future-only benchmark harness with CSV/JSON reports and a quality gate.",
        ],
        0.8,
        1.55,
        8.0,
        3.6,
    )
    add_metric_card(slide, "train windows", "7,992", 9.5, 1.55, BLUE)
    add_metric_card(slide, "benchmark records", "125", 9.5, 3.0, TEAL)
    add_metric_card(slide, "training runtime", "~0.5h", 9.5, 4.45, AMBER)

    slide = new_slide("Experimental Setup", "Fixed 24 FPS continuation task")
    rows = [
        ["Component", "Setting"],
        ["Dataset", "5 held-out Waymo front-camera scenarios"],
        ["Input clips", "20s clips, 10 FPS → 24 FPS via minterpolate"],
        ["Context", "49 frames ≈ 2.0s"],
        ["Generated future", "72 frames = 3.0s"],
        ["Total sequence", "121 frames"],
        ["Seed", "231"],
        ["Training cache", "7,992 train latents, 1,904 val latents"],
    ]
    add_table(slide, rows, 0.75, 1.55, 6.4, 4.25)
    add_bullets(
        slide,
        [
            "All metrics are computed against the real future frames from the same source clip.",
            "Context and future regions are tracked separately so conditioning frames do not inflate the future score.",
        ],
        7.55,
        1.7,
        4.9,
        3.2,
        size=16,
    )

    slide = new_slide("Models Compared")
    rows = [
        ["Model", "Details", "Role"],
        ["Base distilled", "LTX-Video 2B 0.9.8 distilled, no LoRA", "Primary baseline"],
        ["Smoke-test LoRA", "rank 16, LR 2e-5, 512-window subset, step 500", "Failed early recipe"],
        ["Raw-best LoRA", "rank 16, LR 5e-6, step 3000", "Best SSIM/PSNR, fails gate"],
        ["Gate-best LoRA", "rank 8, LR 5e-6, step 1500", "Best valid checkpoint"],
        ["Checkpoint sweep", "ranks 8/16/32 × 8 checkpoints", "Selection study"],
    ]
    add_table(slide, rows, 0.75, 1.55, 11.85, 2.05)
    add_bullets(
        slide,
        [
            "This is not yet action-conditioned training.",
            "The first 512-window LoRA collapsed in sharpness and motion.",
            "The corrected run used the complete existing `latents/` cache; no recaching was required.",
            "The final evaluation now uses 120 LoRA videos plus 5 base references.",
            "The quality gate changes model selection: best raw metric is not the best usable checkpoint.",
        ],
        0.9,
        4.0,
        11.4,
        2.5,
        size=16,
    )

    slide = new_slide("Benchmark Design", "PSNR is not enough for generated video")
    add_bullets(
        slide,
        [
            "Reference metrics: future MSE, MAE, PSNR, and global SSIM.",
            "Quality diagnostics: Laplacian sharpness ratio, motion ratio, temporal delta error, boundary continuity, and context-copy leakage.",
            "Quality gate fails checkpoints that gain pixel score by becoming blurry or under-moving.",
        ],
        0.75,
        1.55,
        6.2,
        4.2,
        size=17,
    )
    add_callout(slide, "Main score region: frames 49–120 only", 7.35, 1.75, 4.85, 0.75, TEAL)
    add_callout(slide, "Context frames are diagnostic, not the main win condition", 7.35, 2.85, 4.85, 0.75, AMBER)

    slide = new_slide("Quantitative Results", "Checkpoint selection changes the conclusion")
    add_table(
        slide,
        [
            ["Model", "Future PSNR ↑", "Future SSIM ↑", "Future MSE ↓", "Sharpness ↑", "Motion ↑"],
            [
                "Base distilled",
                f"{as_float(base, 'mean_future_psnr'):.2f}",
                f"{as_float(base, 'mean_future_global_ssim'):.3f}",
                f"{as_float(base, 'mean_future_mse'):.0f}",
                f"{as_float(base, 'mean_sharpness_ratio_generated_over_reference'):.3f}",
                f"{as_float(base, 'mean_motion_ratio_generated_over_reference'):.3f}",
            ],
            [
                "Best valid: r8 step 1500",
                f"{float(best_gate['future_psnr']):.2f}",
                f"{float(best_gate['future_ssim']):.3f}",
                f"{float(best_gate['future_mse']):.0f}",
                f"{float(best_gate['sharpness_ratio']):.3f}",
                f"{float(best_gate['motion_ratio']):.3f}",
            ],
            [
                "Raw best: r16 step 3000",
                f"{float(raw_best['future_psnr']):.2f}",
                f"{float(raw_best['future_ssim']):.3f}",
                f"{float(raw_best['future_mse']):.0f}",
                f"{float(raw_best['sharpness_ratio']):.3f}",
                f"{float(raw_best['motion_ratio']):.3f}",
            ],
        ],
        0.55,
        1.35,
        12.2,
        1.75,
    )
    slide.shapes.add_picture(str(charts["reference"]), Inches(0.75), Inches(3.25), width=Inches(5.8))
    slide.shapes.add_picture(str(charts["quality"]), Inches(6.85), Inches(3.25), width=Inches(5.8))

    slide = new_slide("Checkpoint Sweep", "Ranks 8/16/32 across all saved checkpoints")
    add_table(
        slide,
        [
            ["Selection", "Rank / Step", "Gate", "PSNR ↑", "SSIM ↑", "Sharp. Ret.", "Motion Ret."],
            [
                "Best valid",
                fmt_rank_step(best_gate),
                "Pass",
                f"{float(best_gate['future_psnr']):.2f}",
                f"{float(best_gate['future_ssim']):.3f}",
                f"{float(best_gate['sharpness_retention_vs_base']):.1%}",
                f"{float(best_gate['motion_retention_vs_base']):.1%}",
            ],
            [
                "Raw SSIM best",
                fmt_rank_step(raw_best),
                "Fail",
                f"{float(raw_best['future_psnr']):.2f}",
                f"{float(raw_best['future_ssim']):.3f}",
                f"{float(raw_best['sharpness_retention_vs_base']):.1%}",
                f"{float(raw_best['motion_retention_vs_base']):.1%}",
            ],
            [
                "Best r16 pass",
                fmt_rank_step(best_by_rank[16]),
                "Pass",
                f"{float(best_by_rank[16]['future_psnr']):.2f}",
                f"{float(best_by_rank[16]['future_ssim']):.3f}",
                f"{float(best_by_rank[16]['sharpness_retention_vs_base']):.1%}",
                f"{float(best_by_rank[16]['motion_retention_vs_base']):.1%}",
            ],
        ],
        0.65,
        1.28,
        12.1,
        1.6,
    )
    slide.shapes.add_picture(str(checkpoint_chart), Inches(1.0), Inches(3.05), width=Inches(11.2))
    add_callout(
        slide,
        "Selected checkpoint: r8 step 1500. Later training is not monotonically better; r16 step 3000 improves SSIM but fails sharpness.",
        1.15,
        6.55,
        11.0,
        0.5,
        AMBER,
    )

    slide = new_slide("Quality Gate Result", "The gate prevents choosing a blurry checkpoint by SSIM alone")
    slide.shapes.add_picture(str(charts["gate"]), Inches(0.85), Inches(1.5), width=Inches(6.4))
    add_bullets(
        slide,
        [
            f"Raw-best r16 step 3000 sharpness retention: {gate['sharpness_retention_lora_over_base']:.1%}; gate requires 80%.",
            f"Raw-best r16 step 3000 future SSIM delta: {gate['future_global_ssim_delta_lora_minus_base']:.4f}.",
            f"Best valid r8 step 1500 future SSIM: {float(best_gate['future_ssim']):.3f}, sharpness retention: {float(best_gate['sharpness_retention_vs_base']):.1%}.",
            "Interpretation: r16 step 3000 is useful evidence that the model learns the domain, but r8 step 1500 is safer to inspect visually.",
        ],
        7.6,
        1.65,
        4.7,
        4.5,
        size=17,
    )
    add_callout(slide, "Best valid checkpoint: r8 step 1500", 7.75, 5.65, 4.2, 0.62, TEAL)

    slide = new_slide("Qualitative Snapshot", "One representative future frame from the 24 FPS evaluation")
    x_positions = [0.8, 4.75, 8.7]
    for (label, image), x in zip(qual_frames, x_positions):
        slide.shapes.add_picture(str(image), Inches(x), Inches(1.65), width=Inches(3.55), height=Inches(3.55))
        box = slide.shapes.add_textbox(Inches(x), Inches(5.3), Inches(3.55), Inches(0.4))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = INK
    add_callout(
        slide,
        "Visual finding: full-data LoRA is much less collapsed than the smoke test, but still slightly over-smooths detail.",
        1.15,
        6.1,
        11.0,
        0.55,
        AMBER,
    )

    slide = new_slide("Analysis: What The Results Mean")
    add_bullets(
        slide,
        [
            "Using the full dataset and LR=5e-6 substantially improves over the 512-window LR=2e-5 smoke test.",
            "r16 step 3000 gives the best raw SSIM/PSNR, but it crosses into over-smoothing.",
            "r8 step 1500 is the best checkpoint that passes SSIM, sharpness, and motion gates together.",
            "Early checkpoints often stay close to base; mid/late checkpoints can trade detail for pixel similarity.",
            "The 49-frame context may still be too short for driving dynamics, and the training objective still needs auditing before action conditioning.",
            "Changing LoRA rank alone is not enough; checkpoint selection is as important as rank selection.",
        ],
        0.8,
        1.45,
        11.6,
        4.9,
        size=17,
    )

    slide = new_slide("What Is Working vs Not Working")
    add_callout(slide, "Working", 0.9, 1.45, 5.2, 0.55, TEAL)
    add_bullets(
        slide,
        [
            "Reproducible Modal inference and local artifact collection.",
            "Fixed-seed base vs LoRA comparisons.",
            "Future-only evaluation harness.",
            "Quality gate catches blur that PSNR misses.",
            "Full-data LoRA improves over the bad 512-window checkpoint.",
            "The rank sweep gives a controlled capacity signal without changing the dataset or frame setup.",
            "The checkpoint sweep found a usable gate-passing checkpoint: r8 step 1500.",
        ],
        0.95,
        2.25,
        5.2,
        3.4,
        size=15,
    )
    add_callout(slide, "Not Working Yet", 7.05, 1.45, 5.2, 0.55, RED)
    add_bullets(
        slide,
        [
            "Full-data LoRA still misses the sharpness gate.",
            "Raw-best r16 step 3000 still misses the sharpness gate.",
            "Only 5 clips have been benchmarked so far.",
            "Context is still only 49 frames, around 2 seconds.",
            "Action conditioning has not been reintroduced yet.",
        ],
        7.1,
        2.25,
        5.2,
        3.4,
        size=15,
    )

    slide = new_slide("Limitations")
    add_bullets(
        slide,
        [
            "Only 5 scenarios in this preliminary benchmark run.",
            "The main full-data benchmark still uses only 5 held-out local scenarios.",
            "We have not yet added LPIPS, DINO/CLIP distance, FVD, or RAFT flow-warp to this harness.",
            "Global SSIM is a lightweight proxy, not a complete perceptual video metric.",
            "Current results evaluate visual continuation only; action-conditioned generation is the next stage.",
        ],
        0.85,
        1.55,
        11.6,
        4.7,
        size=18,
    )

    slide = new_slide("Next Steps")
    add_bullets(
        slide,
        [
            "Visually inspect r8 step 1500 against base and r16 step 3000 on the same five scenarios.",
            "Expand the checkpoint sweep validation set from 5 clips to 50-200 held-out clips.",
            "Increase context to 81 or 121 frames while keeping a 72-frame future.",
            "Add LPIPS/DINO/CLIP and flow-based temporal metrics, then reintroduce action conditioning.",
            "Use r8 step 1500 as the current visual-domain baseline unless larger validation contradicts it.",
        ],
        0.85,
        1.45,
        11.55,
        5.0,
        size=17,
    )

    slide = new_slide("Current Conclusion")
    add_callout(slide, "Best valid checkpoint is r8 step 1500; raw-best r16 step 3000 is sharper evidence of learning but too blurry by gate.", 1.0, 1.7, 11.35, 0.75, BLUE)
    add_bullets(
        slide,
        [
            "The full checkpoint sweep changes the selection from final-step LoRA to an intermediate checkpoint.",
            "r8 step 1500 improves future SSIM over base while retaining 89.2% sharpness and 96.5% motion.",
            "Next step is larger validation plus longer context before scaling to action-conditioned training.",
        ],
        1.0,
        3.0,
        11.0,
        2.6,
        size=19,
    )

    for idx, slide in enumerate(slides, start=1):
        add_footer(slide, idx)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    build_deck()
